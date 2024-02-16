import io
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from urllib.request import urlopen
from io import BytesIO
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import requests
import os


class PresentationGenerator:
    def __init__(self, unsplash_secret_key=None):
        self.unsplash_secret_key = unsplash_secret_key

    def search_unsplash_image(self, query):
        unsplash_url = f"https://api.unsplash.com/photos/random/?query={query}&count=1&client_id={self.unsplash_secret_key}"
        response = requests.get(unsplash_url)
        if response.status_code == 200:
            data = response.json()
            if data:
                return data[0]['urls']['regular']
        return None

    def download_image(self, url):
        response = requests.get(url)
        if response.status_code == 200:
            return response.content
        return None

    def add_image_to_slide_with_search(self, slide, image_name, left, top, width, height):
        formatted_img_name = image_name.split('.')[0].replace('_', ' ')
        print(f"Query param for image - {formatted_img_name}")
        if "none" in formatted_img_name.lower():
            print(f"Corresponding image not found for '{image_name}'. Hence skipping...")
        else:
            image_url = self.search_unsplash_image(formatted_img_name)
            if image_url is None:
                print(f"No image found with name- for '{formatted_img_name}' on Unsplash.")
            else:
                image_data = self.download_image(image_url)
                slide.shapes.add_picture(io.BytesIO(image_data), left, top, width, height)

    def add_text_to_slide(self, slide, text, left, top, width, height, font_size=Pt(18), bold=False, italic=False,
                          color=None, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        for paragraph_text in text.split('\n'):
            p = text_frame.add_paragraph()
            p.text = paragraph_text
            p.font.size = font_size
            p.font.bold = bold
            p.font.italic = italic
            if color and isinstance(color, RGBColor):
                for run in p.runs:
                    run.font.color.rgb = color
            p.alignment = alignment

    def add_multiline_text_to_slide(self, slide, text, left, top, width, height, font_size=Pt(18), bold=False,
                                    italic=False, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT):
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        bullet = '•'
        for sentence in text.split('. '):
            p = text_frame.add_paragraph()
            p.text = f'{bullet} {sentence}.'
            p.font.size = font_size
            p.font.bold = bold
            p.font.italic = italic
            p.alignment = alignment

    def add_image_to_slide(self, slide, image_url, left, top, width, height):
        try:
            image_data = BytesIO(urlopen(image_url).read())
            slide.shapes.add_picture(image_data, left, top, width, height)
        except Exception as e:
            print(f"Error adding image to slide: {e}")

    def is_valid_url(self, url):
        try:
            urlopen(url)
            return True
        except Exception:
            return False

    def create_presentation(self, slide_data, output_file_name):
        presentation = Presentation()
        for index, slide_info in enumerate(slide_data.get('slides', []), start=1):
            blank_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text_frame.text == "Click to Add Title":
                        sp = shape
                        sp.clear()
            title_text = slide_info.get('title', f"Slide {index}")
            self.add_text_to_slide(slide, title_text, Inches(0.5), Inches(0), Inches(9), Inches(1.5),
                                   font_size=Pt(28), bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER)
            intro_text = slide_info.get('intro', '')
            self.add_text_to_slide(slide, intro_text, Inches(0.5), Inches(0.4), Inches(10), Inches(1),
                                   font_size=Pt(14), italic=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER)
            image_name = slide_info.get('image', '')
            if image_name and self.unsplash_secret_key is not None and image_name.lower() not in ['none', 'null', 'nil']:
                self.add_image_to_slide_with_search(slide, image_name, Inches(0.5), Inches(1.5), Inches(4), Inches(5))
            content_text = slide_info.get('content', '')
            self.add_multiline_text_to_slide(slide, content_text, Inches(5.5), Inches(1.5), Inches(3.5), Inches(5),
                                             font_size=Pt(18), alignment=PP_PARAGRAPH_ALIGNMENT.LEFT)
            takeaway_text = f'{slide_info.get("takeaway", "")}'
            text_color = RGBColor(255, 153, 0)
            self.add_text_to_slide(slide, takeaway_text, Inches(0.5), Inches(6.5), Inches(11), Inches(1.5),
                                   font_size=Pt(14), bold=True, color=text_color, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT)

        presentation.save(output_file_name + ".pptx")
